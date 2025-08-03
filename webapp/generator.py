#!/usr/bin/env python3
"""
PowerPoint Song Generator - Web Version
Adapted from simple_generator.py for web application use.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import re
import math
import os


def parse_songs(file_path):
    """Parse songs from text file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
    except UnicodeDecodeError:
        # Try with different encoding
        with open(file_path, 'r', encoding='latin-1') as file:
            content = file.read()
    
    # Split by song markers (lines starting with #)
    song_sections = re.split(r'\n(?=#)', content)
    songs = []
    
    for section in song_sections:
        if section.strip() and section.startswith('#'):
            lines = section.split('\n')
            
            if lines:
                # Extract title (remove #)
                title = lines[0].replace('#', '').strip()
                
                # Get lyrics (everything after title), preserving empty lines
                lyrics = lines[1:] if len(lines) > 1 else []
                
                if title:  # Only add if we have a title
                    songs.append({
                        'title': title,
                        'lyrics': lyrics
                    })
    
    return songs


def split_lyrics_into_slides(lyrics):
    """Split lyrics into slides at paragraph breaks (empty lines)."""
    slides = []
    current_slide = []
    
    for line in lyrics:
        if line.strip():  # Non-empty line
            current_slide.append(line)
        else:  # Empty line - natural paragraph break
            if current_slide:  # Only create slide if we have content
                slides.append(current_slide.copy())
                current_slide = []
    
    # Add remaining content if any
    if current_slide:
        slides.append(current_slide)
    
    return slides


def create_slide(prs, title, content_lines, slide_number=None, total_slides=None):
    """Create a simple slide with title and content."""
    # Use a simple blank layout to avoid placeholder conflicts
    # This ensures consistent behavior across different templates
    try:
        # Use blank layout (index 6) or last available layout
        if len(prs.slide_layouts) > 6:
            layout = prs.slide_layouts[6]  # Blank layout
        else:
            layout = prs.slide_layouts[-1]  # Last available layout
    except (IndexError, AttributeError):
        # Fallback to first available layout
        layout = prs.slide_layouts[0]
    
    slide = prs.slides.add_slide(layout)
    
    # Always use manual text boxes for consistent positioning
    # This avoids conflicts with template placeholders
    
    # Add title at top-left, with distance from red line
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.6),  # More distance from red line at top
        prs.slide_width - Inches(1.8), Inches(1.0)  # Leave minimal space for narrow slide counter
    )
    title_frame = title_box.text_frame
    title_frame.margin_left = Inches(0.2)
    title_frame.margin_right = Inches(0.2)
    title_frame.vertical_anchor = MSO_ANCHOR.TOP
    title_frame.word_wrap = True
    
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)  # Reduced from 40 to 32
    title_p.font.name = "Calibri"
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    title_p.alignment = PP_ALIGN.LEFT
    
    # Add slide counter (e.g., "1/4") in top-right if provided
    if slide_number is not None and total_slides is not None:
        counter_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(1.3), Inches(0.6),  # Top-right position, narrower
            Inches(1.0), Inches(1.0)  # Narrower width for compact counter
        )
        counter_frame = counter_box.text_frame
        counter_frame.margin_left = Inches(0.1)
        counter_frame.margin_right = Inches(0.1)
        counter_frame.vertical_anchor = MSO_ANCHOR.TOP
        counter_frame.word_wrap = False
        
        counter_p = counter_frame.paragraphs[0]
        counter_p.text = f"{slide_number}/{total_slides}"
        counter_p.font.size = Pt(24)
        counter_p.font.name = "Calibri"
        counter_p.font.bold = True
        counter_p.font.color.rgb = RGBColor(139, 69, 19)  # Brown
        counter_p.alignment = PP_ALIGN.RIGHT
    
    # Add content below title, positioned much closer
    if content_lines:
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.4),  # Much closer to title (reduced from 1.8 to 1.4)
            prs.slide_width - Inches(1.0), prs.slide_height - Inches(1.9)  # Adjusted height accordingly
        )
        content_frame = content_box.text_frame
        content_frame.margin_left = Inches(0.2)  # Match title margin
        content_frame.margin_right = Inches(0.2)
        content_frame.vertical_anchor = MSO_ANCHOR.TOP
        content_frame.word_wrap = True
        
        # Add each line with left alignment
        for i, line in enumerate(content_lines):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = line
            p.font.size = Pt(28)  # Reduced from 32 to 28
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(16)


def create_toc_slides(prs, songs_with_slides, songs_per_toc_slide=20):
    """Create Table of Contents slides with clickable links to songs."""
    if not songs_with_slides:
        return []
    
    toc_slides = []
    total_songs = len(songs_with_slides)
    total_toc_slides = math.ceil(total_songs / songs_per_toc_slide)
    
    for toc_page in range(total_toc_slides):
        # Use same layout as song slides
        try:
            if len(prs.slide_layouts) > 6:
                layout = prs.slide_layouts[6]  # Blank layout
            else:
                layout = prs.slide_layouts[-1]  # Last available layout
        except (IndexError, AttributeError):
            layout = prs.slide_layouts[0]
        
        slide = prs.slides.add_slide(layout)
        toc_slides.append(slide)
        
        # Add TOC title
        title_text = f"Table of Contents"
        if total_toc_slides > 1:
            title_text += f" ({toc_page + 1}/{total_toc_slides})"
            
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.6),
            prs.slide_width - Inches(1.0), Inches(1.0)
        )
        title_frame = title_box.text_frame
        title_frame.margin_left = Inches(0.2)
        title_frame.margin_right = Inches(0.2)
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_frame.word_wrap = True
        
        title_p = title_frame.paragraphs[0]
        title_p.text = title_text
        title_p.font.size = Pt(32)
        title_p.font.name = "Calibri"
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(0, 0, 0)
        title_p.alignment = PP_ALIGN.LEFT
        
        # Add songs list for this TOC page in 2 columns
        start_idx = toc_page * songs_per_toc_slide
        end_idx = min(start_idx + songs_per_toc_slide, total_songs)
        songs_on_this_page = songs_with_slides[start_idx:end_idx]
        
        # Calculate columns (10 songs per column)
        col1_songs = songs_on_this_page[:10]  # First 10 songs
        col2_songs = songs_on_this_page[10:]  # Remaining songs
        
        # Create left column (first 10 songs)
        left_column_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6),  # Start below title, left side
            Inches(4.5), prs.slide_height - Inches(2.8)  # Half width, leave space for bottom border
        )
        left_frame = left_column_box.text_frame
        left_frame.margin_left = Inches(0.2)
        left_frame.margin_right = Inches(0.2)
        left_frame.vertical_anchor = MSO_ANCHOR.TOP
        left_frame.word_wrap = True
        
        # Add songs to left column
        for i, (song_title, first_slide_index) in enumerate(col1_songs):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            
            # Add song number and title
            song_number = start_idx + i + 1
            p.text = f"{song_number:2d}. {song_title}"
            p.font.size = Pt(20)  # Larger font for better readability in columns
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue for links
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)  # Moderate space between lines
            
            # Add hyperlink to the song's first slide
            if p.runs:
                run = p.runs[0]
                hlink = run.hyperlink
                hlink.address = f"#{first_slide_index + 1}"
        
        # Create right column (next 10 songs) if there are more songs
        if col2_songs:
            right_column_box = slide.shapes.add_textbox(
                Inches(5.2), Inches(1.6),  # Start below title, right side
                Inches(4.3), prs.slide_height - Inches(2.8)  # Half width, leave space for bottom border
            )
            right_frame = right_column_box.text_frame
            right_frame.margin_left = Inches(0.2)
            right_frame.margin_right = Inches(0.2)
            right_frame.vertical_anchor = MSO_ANCHOR.TOP
            right_frame.word_wrap = True
            
            # Add songs to right column
            for i, (song_title, first_slide_index) in enumerate(col2_songs):
                if i == 0:
                    p = right_frame.paragraphs[0]
                else:
                    p = right_frame.add_paragraph()
                
                # Add song number and title
                song_number = start_idx + len(col1_songs) + i + 1
                p.text = f"{song_number:2d}. {song_title}"
                p.font.size = Pt(20)  # Larger font for better readability in columns
                p.font.name = "Calibri"
                p.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue for links
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(6)  # Moderate space between lines
                
                # Add hyperlink to the song's first slide
                if p.runs:
                    run = p.runs[0]
                    hlink = run.hyperlink
                    hlink.address = f"#{first_slide_index + 1}"
    
    return toc_slides


def generate_presentation(song_file_path, output_path, template_file_path=None, generate_toc=False):
    """
    Generate PowerPoint presentation from song file.
    
    Args:
        song_file_path: Path to the song text file
        output_path: Path where to save the generated PowerPoint
        template_file_path: Optional path to PowerPoint template
        generate_toc: Whether to generate table of contents
    
    Returns:
        tuple: (success: bool, message: str, slide_count: int)
    """
    try:
        # Parse songs
        songs = parse_songs(song_file_path)
        if not songs:
            return False, "No songs found in the file. Make sure song titles start with #", 0
        
        # Create presentation
        if template_file_path and os.path.exists(template_file_path):
            prs = Presentation(template_file_path)
        else:
            prs = Presentation()
        
        total_slides = 0
        songs_with_slide_positions = []  # Track song titles and their first slide positions
        
        # If TOC is requested, we need to calculate TOC slides first to adjust positions
        toc_slides_count = 0
        if generate_toc:
            # First pass: collect song titles
            temp_songs = []
            for song in songs:
                temp_songs.append(song['title'])
            # Calculate how many TOC slides we'll need
            toc_slides_count = math.ceil(len(temp_songs) / 20)  # 20 songs per TOC slide (2 columns x 10 rows)
        
        for song in songs:
            title = song['title']
            lyrics = song['lyrics']
            
            # Split lyrics into slides
            lyric_slides = split_lyrics_into_slides(lyrics)
            
            # Record the first slide position for this song (adjust for TOC at beginning)
            first_slide_position = total_slides + toc_slides_count
            songs_with_slide_positions.append((title, first_slide_position))
            
            # Create slides for this song with numbering
            total_song_slides = len(lyric_slides)
            for slide_index, slide_content in enumerate(lyric_slides):
                create_slide(prs, title, slide_content, slide_index + 1, total_song_slides)
                total_slides += 1
        
        # Generate Table of Contents if requested - create at beginning
        if generate_toc and songs_with_slide_positions:
            # Create new presentation with TOC first
            if template_file_path and os.path.exists(template_file_path):
                final_prs = Presentation(template_file_path)
            else:
                final_prs = Presentation()
            
            # Remove default slides
            slide_count = len(final_prs.slides)
            for i in range(slide_count):
                rId = final_prs.slides._sldIdLst[0].rId
                final_prs.part.drop_rel(rId)
                del final_prs.slides._sldIdLst[0]
            
            # Create TOC slides first
            create_toc_slides(final_prs, songs_with_slide_positions)
            
            # Add song slides after TOC - recreate them with slide numbering
            for song in songs:
                title = song['title']
                lyrics = song['lyrics']
                
                # Split lyrics into slides
                lyric_slides = split_lyrics_into_slides(lyrics)
                total_song_slides = len(lyric_slides)
                
                # Create slides for this song with numbering
                for slide_index, slide_content in enumerate(lyric_slides):
                    create_slide(final_prs, title, slide_content, slide_index + 1, total_song_slides)
            
            prs = final_prs
            total_slides += toc_slides_count
        
        # Save presentation
        prs.save(output_path)
        
        # Return success
        success_message = f"Generated {total_slides} slides from {len(songs)} songs"
        if generate_toc and toc_slides_count > 0:
            success_message += f" + {toc_slides_count} TOC slides"
        
        return True, success_message, total_slides
        
    except FileNotFoundError as e:
        return False, f"File not found: {str(e)}", 0
    except Exception as e:
        return False, f"Error generating presentation: {str(e)}", 0


if __name__ == "__main__":
    # Test the generator
    import sys
    if len(sys.argv) < 3:
        print("Usage: python generator.py <song_file> <output_file> [template_file] [--toc]")
        sys.exit(1)
    
    song_file = sys.argv[1]
    output_file = sys.argv[2]
    template_file = sys.argv[3] if len(sys.argv) > 3 and not sys.argv[3].startswith('--') else None
    generate_toc = '--toc' in sys.argv
    
    success, message, slide_count = generate_presentation(song_file, output_file, template_file, generate_toc)
    
    if success:
        print(f"✅ {message}")
    else:
        print(f"❌ {message}")
        sys.exit(1)